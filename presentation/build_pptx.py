#!/usr/bin/env python3
"""Build the editable PowerPoint version of the Harris-Nelson Family Reunion
presentation. The .pptx opens in PowerPoint and imports into Canva, so names,
captions, and numbers can be edited later.

Usage:  python3 build_pptx.py
Output: Harris-Nelson-Reunion-2026.pptx (same folder)
Photos: pulled from ../website/photos/ (family/ and 2024/)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
PHOTOS = os.path.join(HERE, "..", "website", "photos")

GREEN = RGBColor(0x1B, 0x5E, 0x20)
GREEN_MID = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
RED = RGBColor(0xB7, 0x1C, 0x1C)
CREAM = RGBColor(0xFA, 0xF6, 0xEE)
INK = RGBColor(0x26, 0x22, 0x1B)
MUTED = RGBColor(0x6B, 0x63, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xF5, 0xD5, 0x8A)
DARK = RGBColor(0x21, 0x1D, 0x17)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = 13.333, 7.5


def slide(bg=CREAM):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def text(s, x, y, w, h, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         space_after=6, font="Georgia"):
    """runs: string, or list of paragraphs; a paragraph may be (text, dict)."""
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        opts = {}
        if isinstance(para, tuple):
            para, opts = para
        p.text = para
        p.alignment = opts.get("align", align)
        p.space_after = Pt(opts.get("space_after", space_after))
        f = p.runs[0].font if p.runs else p.font
        f.size = Pt(opts.get("size", size))
        f.bold = opts.get("bold", bold)
        f.color.rgb = opts.get("color", color)
        f.name = opts.get("font", font)
    return box


def kicker(s, t, y=0.35, color=ORANGE):
    text(s, 0.7, y, W - 1.4, 0.4, t.upper(), size=13, color=color, bold=True,
         font="Verdana")


def title(s, t, y=0.75, size=34, color=GREEN):
    text(s, 0.7, y, W - 1.4, 0.9, t, size=size, color=color, bold=True)


def bullets(s, x, y, w, items, size=15, head=None, head_color=GREEN_MID, h=4.5):
    paras = []
    if head:
        paras.append((head, {"size": size + 2, "bold": True, "color": head_color,
                             "space_after": 8}))
    for it in items:
        paras.append(("• " + it, {"size": size}))
    text(s, x, y, w, h, paras)


def photo_row(s, files, y, ph, caption="[Add name]", x0=0.7, gap=0.25):
    """Place photos in a row with editable caption boxes underneath."""
    files = [f for f in files if os.path.exists(f)]
    if not files:
        return
    n = len(files)
    pw = (W - 2 * x0 - gap * (n - 1)) / n
    for i, f in enumerate(files):
        x = x0 + i * (pw + gap)
        pic = s.shapes.add_picture(f, Inches(x), Inches(y), height=Inches(ph))
        # center horizontally in its cell, clamp width
        if pic.width > Inches(pw):
            pic.width = Inches(pw)
            pic.left = Inches(x)
        else:
            pic.left = Inches(x + (pw - pic.width / 914400) / 2)
        text(s, x, y + ph + 0.05, pw, 0.35, caption, size=12, color=MUTED,
             align=PP_ALIGN.CENTER, font="Verdana")


def footer(s, t, color=MUTED):
    text(s, 0.7, H - 0.62, W - 1.4, 0.5, t, size=11, color=color, font="Verdana")


P = lambda *parts: os.path.join(PHOTOS, *parts)

# ---------------------------------------------------------------- 1 · TITLE
s = slide()
text(s, 0.7, 1.15, W - 1.4, 0.5, "TOGETHER AGAIN! · “HERE WE GO AGAIN… KICK’N IT WITH THE KINS!”",
     size=14, color=ORANGE, bold=True, align=PP_ALIGN.CENTER, font="Verdana")
text(s, 0.7, 1.8, W - 1.4, 2.2, "Harris-Nelson\nFamily Reunion 2026", size=54,
     color=GREEN, bold=True, align=PP_ALIGN.CENTER)
text(s, 0.7, 4.2, W - 1.4, 0.6, "Our roots run deep — from Mississippi to Cleveland and beyond…",
     size=20, color=MUTED, align=PP_ALIGN.CENTER)
text(s, 0.7, 5.0, W - 1.4, 0.5, "Cleveland, Ohio · July 17–19, 2026 · hnfamilyreunion.com",
     size=14, color=MUTED, align=PP_ALIGN.CENTER, font="Verdana")

# ---------------------------------------------------------------- 2 · AGENDA
s = slide()
kicker(s, "Today's Program"); title(s, "What We'll Cover")
left = ["Welcome & The Weekend at a Glance", "Our Family History & Trivia",
        "Family Yearbook", "Reunion 2024 Memories — Cleveland",
        "Acknowledgments & Accomplishments", "Family Superlatives — Cast Your Vote!",
        "Tribute to Loved Ones We've Lost", "Family Business: Financial Report",
        "Officers, Robert's Rules & Our Constitution"]
right = ["Board Recognition, Committees & Sign-Ups", "Where We Meet Next — Location Vote",
         "Our New Website: hnfamilyreunion.com", "Getting Our Land Back — Where Is My Land",
         "Scholarship Fund for Our Descendants", "Helping Hands — Family Hardship Fund",
         "The Family Foundation — Our 501(c)(3)",
         "The Family Tree Project — Add Yourself!", "Closing & Next Steps"]
text(s, 0.9, 1.8, 5.9, 5.2, [(f"{i+1}.  {t}", {"size": 16, "space_after": 10}) for i, t in enumerate(left)])
text(s, 7.0, 1.8, 5.9, 5.2, [(f"{i+10}.  {t}", {"size": 16, "space_after": 10}) for i, t in enumerate(right)])

# ------------------------------------------------------- 3 · WEEKEND SCHEDULE
s = slide()
kicker(s, "July 17–19, 2026 · Cleveland, Ohio"); title(s, "The Weekend at a Glance")
bullets(s, 0.7, 1.8, 4.0, [
    "Andrews Osborne Park (in back), 38575 Lakeshore Blvd, Willoughby OH 44094",
    "Wear any PREVIOUS reunion shirt",
    "Pick up your 2026 shirt & Clay's Park tickets",
    "Courts, pool, playground, disc golf & swings by the lake"],
    head="FRIDAY 7/17 · 5–9 PM · Meet & Greet / Fish Fry", head_color=GREEN_MID)
bullets(s, 4.85, 1.8, 4.0, [
    "Clay's Park, 12951 Patterson St NW, North Lawrence OH 44666",
    "Wear the BEIGE/BURGUNDY 2026 reunion tee",
    "Water & dry activities all day",
    "FAMILY PICTURE 11:30 AM at the pavilion — ALL attendees"],
    head="SATURDAY 7/18 · 10 AM–3 PM · Family Fun", head_color=ORANGE)
bullets(s, 9.0, 1.8, 3.9, [
    "SiteCenters, 3333 Richmond Rd, Beachwood OH",
    "Wear WHITE top & blue-jean bottom",
    "Brunch 11:30 AM · family pictures 1:45 PM",
    "Acknowledgments, next reunion, committee changes & close-out"],
    head="SUNDAY 7/19 · 11 AM–3 PM · Brunch & Farewell", head_color=RED)
footer(s, "Shirts: $10 (2XL–4XL $15) — adult SM–4XL, youth SM–XL, toddler 6M–24M. Pick up Friday; wear Saturday.")

# ---------------------------------------------------------------- 4 · WELCOME
s = slide()
kicker(s, "Welcome"); title(s, "Welcome to the Harris-Nelson Family Reunion")
text(s, 0.7, 1.9, 6.2, 4.5, [
    ("This reunion started when Mr. John Harris and Ms. Judy Bender united a long time ago. Both were born in the late 1800s — he was a farmer, and she was a homemaker.", {"size": 17, "space_after": 12}),
    ("Their union brought forth six children, four boys and two girls: Gussy, Pasty, Josh, A.C., D.B., and Will Harris.", {"size": 17, "space_after": 12}),
    ("“As I call your name, please stand with your family!”", {"size": 19, "bold": True, "color": ORANGE})])
bullets(s, 7.4, 1.9, 5.2, [
    "To honor the people who got us here",
    "To make sure every branch knows every other branch",
    "To pass the story — and the land, and the legacy — to the next generation"],
    head="Why we gather", size=16)

# ---------------------------------------------------------------- 5 · HISTORY
s = slide()
kicker(s, "Our Family History"); title(s, "From One Union, Many Branches")
text(s, 0.7, 1.75, W - 1.4, 5.0, [
    ("1 · John Harris & Judy Bender (born late 1800s) — six children: Gussy, Pasty, Josh, A.C., D.B., and Will Harris.", {"size": 16, "space_after": 12}),
    ("2 · Emma (Bently) Moncrief (1867–1954) & Ed Moncrief had 16 children. Out of these 16, Mandy Moncrief (c. 1900–1986) was born. Mandy married Nelson Price (b. 1894) and blessed us with Lela Mae Price (July 16, 1922 – February 15, 1996).", {"size": 16, "space_after": 12}),
    ("3 · Lela Mae Price married Dennis Bernard “D.B.” Harris — three girls: Mary Nelson, Mildred Lucile Harris (1942–2012), Jessie Mae Harris-Moore. She later married Floyd Nelson and had six more children: James Edward, James Earl (mom liked the name James!), Floyd Nelson Jr, Mary Jane Harris (Gray), Mandy Louise Ellis, and the baby of the bunch, Shirley Harris.", {"size": 16, "space_after": 12}),
    ("4 · Charlie Nelson & Mandy Harris brought forth two boys: Floyd Nelson (the oldest) and Joe Nelson. Joe married Oreatha Ward — four children: Doris Joe Smith, Mandy Mary Thomas, Charles Nelson, and the baby, Naiomi Matthews.", {"size": 16, "space_after": 12})])
footer(s, "This family history was brought to you by Yours Truly — That Aunt \U0001F49B")

# --------------------------------------------------------------- 6 · BRANCHES
s = slide()
kicker(s, "Our Family History"); title(s, "The Nine Branches — Stand With Your Family!")
branches = [("Mary Nelson", 8), ("Mildred Ellis", 6), ("Jessie Moore", 4),
            ("James Edward", 5), ("James Earl", 3), ("Curtis", 4),
            ("Mary Jane Gray", 2), ("Mandy Ellis", 3), ("Shirley Harris", 1)]
for i, (name, n) in enumerate(branches):
    col, row = i % 3, i // 3
    x, y = 0.8 + col * 4.05, 1.9 + row * 1.5
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(3.8), Inches(1.1))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = name; p.runs[0].font.size = Pt(20)
    p.runs[0].font.bold = True; p.runs[0].font.color.rgb = GREEN
    p2 = tf.add_paragraph(); p2.text = f"{n} children"
    p2.runs[0].font.size = Pt(15); p2.runs[0].font.color.rgb = ORANGE
footer(s, "Every branch keeps growing — that's what the Family Tree Project is for!")

# ----------------------------------------------------------------- 7 · TRIVIA
s = slide()
kicker(s, "Let's Play"); title(s, "Family History Trivia \U0001F3C6")
qa = [("Who founded this family, and what did he do for a living?", "John Harris (m. Judy Bender) — a farmer; she was a homemaker"),
      ("How many children did Emma & Ed Moncrief have?", "Sixteen!"),
      ("What is Lela Mae Price's birthday?", "July 16, 1922"),
      ("Which two brothers share a first name — and why?", "James Edward & James Earl — mom just liked the name James!"),
      ("Who is “the baby of the bunch”?", "Shirley Harris"),
      ("Joe Nelson married whom?", "Oreatha Ward — 4 children"),
      ("Which branch has the most children?", "Mary Nelson's, with 8"),
      ("In what year did Mary marry Willie Nelson?", "1957")]
paras = []
for i, (q, a) in enumerate(qa, 1):
    paras.append((f"Q{i}. {q}", {"size": 14, "bold": True, "space_after": 2}))
    paras.append((f"      A: {a}", {"size": 13, "color": ORANGE, "space_after": 8}))
half = len(paras) // 2
text(s, 0.7, 1.75, 6.0, 5.2, paras[:half])
text(s, 6.95, 1.75, 6.0, 5.2, paras[half:])

# --------------------------------------------------------------- 8 · YEARBOOK
s = slide()
kicker(s, "Keepsake"); title(s, "The Family Yearbook \U0001F4D6")
bullets(s, 0.7, 1.9, 6.0, [
    "Portraits, superlatives, branch pages, recipes & the year's biggest moments",
    "Submit photos & captions to harrisnelsonfamilyreunion@gmail.com",
    "Photographer & videographer shoot portraits at reunion events",
    "Recipe/cookbook section — bring your best dish AND the recipe"], size=16)
bullets(s, 7.2, 1.9, 5.4, [
    "Yearbook preview post (X): x.com/tendin2",
    "Yearbook & history working doc: Google Docs (link on the website)",
    "Photos & stories album: the family Canva album"],
    head="Yearbook links", size=15)

# ------------------------------------------------------- 9 · 2024 MEMORIES
s = slide()
kicker(s, "Reunion 2024 · Cleveland, Ohio")
title(s, "“We All We Got… We All We Need!” \U0001F4F8")
photo_row(s, [P("2024", f"{n}.jpg") for n in (1, 15, 45, 105)], 1.9, 1.75,
          caption="[Add caption]")
bullets(s, 0.7, 4.3, 6.0, [
    "All 123 slides play fullscreen on the Photos page of hnfamilyreunion.com",
    "Put it on the big screen during dinner!"],
    head="▶ Watch the full slideshow", size=15, h=2.2)
bullets(s, 7.2, 4.3, 5.4, [
    "Send this weekend's photos to harrisnelsonfamilyreunion@gmail.com",
    "The Photography Committee builds the next slideshow & yearbook from what you send"],
    head="Add to next year's show", size=15, h=2.2)

# ------------------------------------------------- 10 · ACKNOWLEDGMENTS/GRADS
s = slide()
kicker(s, "Celebrate"); title(s, "Acknowledgments — Congratulations, Class of 2026! \U0001F393")
photo_row(s, [P("family", f"grad-{i}.jpg") for i in range(1, 6)], 1.8, 2.6)
text(s, 0.7, 4.95, W - 1.4, 1.8, [
    ("Also celebrating: [new babies] · [weddings & anniversaries] · [promotions, retirements & new businesses] · [military service, awards & honors] · [our elders present today]",
     {"size": 15, "color": MUTED})])
footer(s, "Edit the [Add name] boxes above with each graduate's name and school. Send accomplishments to harrisnelsonfamilyreunion@gmail.com for the yearbook.")

# ------------------------------------------- 10b · QUESTIONNAIRE SPOTLIGHT
s = slide()
kicker(s, "Straight From the Family Questionnaire")
title(s, "You Talked, We Listened \U0001F4DD")
bullets(s, 0.7, 1.8, 4.0, [
    "Miesha — first college degree in her grandmother's family; NuLife Fitness Camp turns 20 (Mar 2026); Juicy Vegan; property owner since 2021",
    "Lovie — homeowner at 24 (May 2026)",
    "Denisha — RN for 10 years, NP in progress",
    "Dennis Jr. — serves special-needs students"],
    head="\U0001F31F Shine on", size=12)
bullets(s, 4.85, 1.8, 4.0, [
    "Nyeri is our resident makeup artist",
    "Lovie reps us in Cincinnati — and Delta Sigma Theta Sorority, Inc.",
    "Mary Jane loves to fish",
    "Miesha researched the census from when Grandma Lela Mae was a baby"],
    head="\U0001F604 Did you know…", size=12, head_color=ORANGE)
bullets(s, 9.0, 1.8, 3.9, [
    "More participants — every branch in the building",
    "More creative ways to showcase our history",
    "Skits from the youth",
    "A destination reunion — 'a 1-hour trip'",
    "Keep the competitions coming!"],
    head="\U0001F5F3 What y'all want next", size=12, head_color=RED)
footer(s, "From the 2026 Family Questionnaire (7 responses so far). Genealogy answers went straight to the Historian for the family tree.")

# ------------------------------------------- 10c · PARTICIPATION STATS
s = slide()
kicker(s, "Real Talk · The Numbers")
title(s, "115 of Us Showed Up. 7 Filled Out the Questionnaire. \U0001F62C")
for i, (amt, lab, c) in enumerate([("≈115", "Family members at Reunion 2026", GREEN),
                                   ("7", "Questionnaire responses", ORANGE),
                                   ("6%", "Our participation rate, family…", RED)]):
    x = 0.8 + i * 4.15
    text(s, x, 1.75, 3.9, 1.3, [(amt, {"size": 34, "bold": True, "color": c, "space_after": 2}),
                                 (lab, {"size": 13, "color": MUTED})])
TRACK = RGBColor(0xEE, 0xE4, 0xD0)
def hbar(x, y, label, value_txt, frac, w_lab=2.55, w_bar=2.5, track=False, fill=GREEN_MID):
    text(s, x, y - 0.04, w_lab, 0.3, label, size=10, color=MUTED, font="Verdana")
    if track:
        bg = s.shapes.add_shape(1, Inches(x + w_lab), Inches(y), Inches(w_bar), Inches(0.22))
        bg.fill.solid(); bg.fill.fore_color.rgb = TRACK; bg.line.fill.background()
    if frac > 0:
        fg = s.shapes.add_shape(1, Inches(x + w_lab), Inches(y), Inches(max(0.06, w_bar * frac)), Inches(0.22))
        fg.fill.solid(); fg.fill.fore_color.rgb = fill; fg.line.fill.background()
    text(s, x + w_lab + w_bar + 0.08, y - 0.04, 0.9, 0.3, value_txt, size=10, color=INK, bold=True, font="Verdana")
text(s, 0.8, 2.95, 5.6, 0.3, "Who showed up vs. who spoke up", size=12, bold=True, color=GREEN)
hbar(0.8, 3.35, "Attended Reunion 2026", "≈115", 1.0)
hbar(0.8, 3.72, "Serving on the board today", "8", 0.07)
hbar(0.8, 4.09, "Filled out the questionnaire", "7", 0.06)
text(s, 7.0, 2.95, 5.6, 0.3, "Seats filled vs. seats we need", size=12, bold=True, color=GREEN)
hbar(7.0, 3.35, "Board seats", "8/12", 0.67, w_lab=2.3, w_bar=2.3, track=True)
hbar(7.0, 3.72, "Committee chairs", "1/11", 0.09, w_lab=2.3, w_bar=2.3, track=True)
hbar(7.0, 4.09, "Committee members", "0/~33", 0.0, w_lab=2.3, w_bar=2.3, track=True)
bullets(s, 0.7, 4.55, 6.0, [
    "The yearbook, tree poster, acknowledgments & trivia were ALL built from those 7 responses",
    "Imagine what we could show with 115 — your history dies untold if you don't tell it"],
    head="Why it matters", size=12, h=1.9)
bullets(s, 7.0, 4.55, 5.6, [
    "Before you leave Sunday: fill out the questionnaire — phone, paper, or grab a grandchild to type",
    "Join a committee — 9 chairs and every membership seat still open",
    "Add yourself to the family tree; send photos to the family email"],
    head="The plea \U0001F64F", size=12, head_color=RED, h=1.9)
footer(s, "Goal for the next reunion: 100+ questionnaire responses. The Membership Chair (voting today!) will be on all of us about it — lovingly.")

# --------------------------------------------------------- 11 · FAMILY ALBUM
s = slide()
kicker(s, "Family Album"); title(s, "Throwbacks & Treasures \U0001F4BE")
photo_row(s, [P("family", f"album-{i}.jpg") for i in range(1, 5)], 1.75, 2.2)
photo_row(s, [P("family", f"album-{i}.jpg") for i in range(5, 9)], 4.45, 2.2)
footer(s, "Edit the [Add name] boxes with who's who — these photos also live on the website's Photos page.")

# ---------------------------------------------------------- 12 · SUPERLATIVES
s = slide()
kicker(s, "Vote All Weekend"); title(s, "Family Superlatives \U0001F3C5 — Cast Your Vote!")
bullets(s, 0.7, 1.8, 4.0, ["Best Dressed", "Life of the Party", "Best Dancer",
                            "Loudest Laugh", "Most Likely to Be Late"],
        head="\U0001F602 The Fun Ones", size=15)
bullets(s, 4.85, 1.8, 4.0, ["Best Cook / Best Dish", "Grill Master",
                             "Best Card/Domino Player", "Spades Champion's Table",
                             "MVP of the Family Games"],
        head="\U0001F357 The Skills", size=15, head_color=ORANGE)
bullets(s, 9.0, 1.8, 3.9, ["Family Glue Award", "Traveled the Farthest",
                            "Eldest & Youngest Present", "Biggest Branch in Attendance",
                            "Legacy Keeper (history & stories)"],
        head="\U0001F49B The Heart", size=15, head_color=RED)
footer(s, "Nominations open Friday at check-in · voting closes Saturday 6 PM · one ballot per person · Hospitality counts · winners crowned at Sunday brunch.")

# --------------------------------------------------------------- 13 · TRIBUTE
s = slide(bg=DARK)
kicker(s, "In Loving Memory", color=GOLD)
title(s, "A Tribute to Those We've Lost \U0001F54A️", color=GOLD)
text(s, 0.7, 1.9, W - 1.4, 1.2,
     "They are not gone — they are the roots beneath everything we do today. Please stand with us for a moment of silence.",
     size=18, color=RGBColor(0xCA, 0xBF, 0xA8))
text(s, 0.7, 3.2, 6.0, 2.6, [
    ("We remember", {"size": 17, "bold": True, "color": GOLD, "space_after": 8}),
    ("[Add the names of loved ones lost, with branch and years]", {"size": 15, "color": WHITE})])
text(s, 7.0, 3.2, 5.6, 2.6, [
    ("Candle-lighting", {"size": 17, "bold": True, "color": GOLD, "space_after": 8}),
    ("One representative from each of the nine branches lights a candle as their branch's names are read aloud.", {"size": 15, "color": WHITE})])
footer(s, "A memorial page in the yearbook and on the website keeps their stories alive.", color=RGBColor(0x94, 0x8A, 0x75))

# ------------------------------------------------------ 14 · BUSINESS DIVIDER
s = slide(bg=GREEN)
text(s, 0.7, 2.6, W - 1.4, 1.2, "Family Business", size=52, color=WHITE,
     bold=True, align=PP_ALIGN.CENTER)
text(s, 0.7, 4.0, W - 1.4, 0.7,
     "Finances · Governance · Committees · Website · Land · Scholarships · Foundation",
     size=17, color=RGBColor(0xD9, 0xE7, 0xD9), align=PP_ALIGN.CENTER)

# ------------------------------------------------------------- 15 · FINANCIAL
s = slide()
kicker(s, "Family Business · Financial Report")
title(s, "Where the Money Stands \U0001F4B5  (as of July 13, 2026)")
for i, (amt, lab, c) in enumerate([("$5,540", "Dues collected · 115 attending", GREEN),
                                   ("$1,675", "Dues still outstanding", ORANGE),
                                   ("$410", "Projected shortfall — let's close it!", RED)]):
    x = 0.8 + i * 4.15
    text(s, x, 1.75, 3.9, 1.3, [(amt, {"size": 30, "bold": True, "color": c, "space_after": 2}),
                                 (lab, {"size": 13, "color": MUTED})])
rows = [("Clay's Park pavilion", "$700", "PAID 3/13/26"),
        ("Friday pavilion (Andrews Osborne)", "$65", "PAID 3/1/26"),
        ("Clay's Park tickets — 93 @ $20", "$1,860", "Due"),
        ("Sunday banquet space", "$200", "Due"),
        ("Sunday brunch catering", "$2,600", "$1,734 paid · $866.75 due"),
        ("Estimated food — Fri / Sat", "$2,200", "Due"),
        ("TOTAL", "$7,625", "$2,499 paid · $5,127 remaining")]
tbl = s.shapes.add_table(len(rows) + 1, 3, Inches(0.8), Inches(3.15),
                         Inches(11.7), Inches(3.0)).table
tbl.columns[0].width = Inches(5.4); tbl.columns[1].width = Inches(2.0)
tbl.columns[2].width = Inches(4.3)
for j, htxt in enumerate(["Expense", "Amount", "Status"]):
    c = tbl.cell(0, j); c.text = htxt
    c.fill.solid(); c.fill.fore_color.rgb = GREEN
    r = c.text_frame.paragraphs[0].runs[0].font
    r.size = Pt(12); r.bold = True; r.color.rgb = WHITE
for i, row in enumerate(rows, 1):
    for j, v in enumerate(row):
        c = tbl.cell(i, j); c.text = v
        f = c.text_frame.paragraphs[0].runs[0].font
        f.size = Pt(12); f.color.rgb = INK
        f.bold = (row[0] == "TOTAL")
footer(s, "Dues total expected: $7,215. Dues: $125 single (or w/ 1 minor) · $175 single + guest/college student · $225 family of 3+. See the Treasurer for your balance.")

# -------------------------------------------------------------- 16 · OFFICERS
s = slide()
kicker(s, "Family Business · Governance"); title(s, "How We'll Govern Going Forward \U0001F3DB")
officers = [("President", "Miesha", "Presides, sets agenda, speaks for the family"),
            ("Vice President", "Aunt Shirley", "Presides in absence; oversees committee chairs"),
            ("Secretary", "Marcelette", "Minutes; TAKES ATTENDANCE; records, correspondence & notices"),
            ("Treasurer", "Jasmine", "Holds funds, pays approved expenses, reports"),
            ("Financial Secretary", "Aunt Vanessa", "Records all money in, issues receipts, tracks dues"),
            ("Historian", "Shone", "History, archives, yearbook, family tree"),
            ("Hospitality Chair", "Loretta", "Leads the COURTESY COMMITTEE: welcome, elders' care, tables & ushers, birthdays & accomplishments"),
            ("Sergeant-at-Arms", "Joseph", "Keeps order; GUARDS THE DOOR (assisted by the Doorkeeper); safety"),
            ("Membership Chair", "[Open — vote today]", "Membership roll & recruitment; matches volunteers to committees"),
            ("Parliamentarian", "[Open — vote today]", "Advises on Robert's Rules; rules on points of order; CHAIRS THE PROTOCOL COMMITTEE"),
            ("Doorkeeper", "[Open — vote today]", "Assists the Sergeant-at-Arms: door check-in, controls entry during votes, carries messages"),
            ("Meditation Chair", "[Open — vote today]", "Leads prayers & prayer requests; delivers the condolences report; memorials with Services")]
tbl = s.shapes.add_table(len(officers) + 1, 3, Inches(0.8), Inches(1.55),
                         Inches(11.7), Inches(5.3)).table
tbl.columns[0].width = Inches(2.9); tbl.columns[1].width = Inches(2.4)
tbl.columns[2].width = Inches(6.4)
for j, htxt in enumerate(["Office", "Officer", "Core duty"]):
    c = tbl.cell(0, j); c.text = htxt
    c.fill.solid(); c.fill.fore_color.rgb = GREEN
    r = c.text_frame.paragraphs[0].runs[0].font
    r.size = Pt(12); r.bold = True; r.color.rgb = WHITE
for i, row in enumerate(officers, 1):
    for j, v in enumerate(row):
        c = tbl.cell(i, j); c.text = v
        f = c.text_frame.paragraphs[0].runs[0].font
        f.size = Pt(10); f.color.rgb = INK; f.bold = (j == 1)
footer(s, "Our written Constitution & Bylaws (Robert's Rules as parliamentary authority) is ready for adoption at this meeting.")

# ---------------------------------------------------- 16b · RECOGNITION
s = slide()
kicker(s, "Family Business · Recognition")
title(s, "Give Our Board & Committee Members Their Flowers \U0001F44F")
bullets(s, 0.7, 1.8, 4.0, [
    "Miesha — President", "Aunt Shirley — Vice President", "Marcelette — Secretary",
    "Jasmine — Treasurer", "Aunt Vanessa — Financial Secretary", "Shone — Historian",
    "Loretta — Hospitality", "Joseph — Sergeant-at-Arms"],
    head="Our current board — thank you!", size=13)
bullets(s, 4.85, 1.8, 4.0, [
    "[Add the names of everyone currently serving on a committee — they get recognized here]"],
    head="Our committee members", size=13, head_color=ORANGE)
bullets(s, 9.0, 1.8, 3.9, [
    "Four new board seats, open for votes today:",
    "MEMBERSHIP CHAIR — roll & recruitment", "PARLIAMENTARIAN — Robert's Rules + Protocol Committee",
    "DOORKEEPER — assists the Sergeant-at-Arms", "MEDITATION CHAIR — prayers & condolences report",
    "COMMITTEE MEMBERS — every committee needs 2–5 members to assist its chair; no chair should carry a committee alone"],
    head="We need YOU", size=13, head_color=RED)
footer(s, "The Membership Chair keeps the family roll, welcomes new members, and matches volunteers to committees. The Secretary is taking nominations right now.")

# -------------------------------------------------------- 17 · ROBERT'S RULES
s = slide()
kicker(s, "Family Business · Governance")
title(s, "Robert's Rules of Order + Our Constitution \U0001F4DC")
bullets(s, 0.7, 1.8, 6.0, [
    "Quorum — minimum members present to conduct business",
    "Motion — “I move that we…” puts an idea on the floor",
    "Second — one member seconds to open discussion",
    "Discussion — the chair recognizes speakers in turn",
    "Vote — majority carries; the Secretary records it"],
    head="Meetings stay fair, short & drama-free", size=15)
bullets(s, 7.0, 1.8, 5.6, [
    "Propose → “I move that…”",
    "Support → “Second!”",
    "Change wording → “I move to amend…”",
    "End debate → “I move the previous question”",
    "Wrap up → “I move to adjourn”"],
    head="Cheat sheet — what do I say?", size=15, head_color=ORANGE)
footer(s, "Motion on today's agenda: adopt the Harris-Nelson Constitution & Bylaws — membership, officers & terms, quorum, dues, committees, amendments.")

# ------------------------------------------------- 17b · OPEN CHAIRS
s = slide()
kicker(s, "Family Business · Open Seats")
title(s, "Nine Committee Chairs Still Need a Name Next to Them \U0001F4CB")
chairs = [
    ("\U0001F3C8 Sports / Games", "tournaments, kids' games, family Olympics, spades & dominoes brackets"),
    ("\U0001F3B6 Music", "DJ & playlists, line-dance sets, old-school hour, banquet music"),
    ("\U0001F37D Food", "caterers, menus, the cookout & potluck, family cookbook"),
    ("\U0001F9BA Safety", "first aid, headcounts, weather plans — with the Sergeant-at-Arms"),
    ("\U0001F64F Services", "Sunday worship, memorial tribute, elder care — with the Meditation Chair"),
    ("\U0001F4F8 Photography", "photographer & videographer, branch portraits, yearbook photos"),
    ("\U0001F49D Philanthropy & Fundraising", "fundraisers for all three funds; scholarship & hardship review"),
    ("\U0001F5D3 Planning", "dates, venues, decade planning, budgets, registration"),
    ("\U0001F4BB Technology", "hnfamilyreunion.com, online payments & shirts, tree tools")]
for i, (cname, duty) in enumerate(chairs):
    col, row = i % 3, i // 3
    x, y = 0.8 + col * 4.1, 1.75 + row * 1.5
    text(s, x, y, 3.9, 1.4, [(cname + " Chair — OPEN", {"size": 14, "bold": True, "color": ORANGE, "space_after": 1}),
                              (duty, {"size": 11, "color": MUTED})])
footer(s, "Courtesy (Loretta) and Protocol (Parliamentarian) are covered. Chairs appointed by the President, confirmed by member vote — raise your hand or check 'I want to chair' on the sign-up form.")

# ------------------------------------------------------------ 18 · COMMITTEES
s = slide()
kicker(s, "Family Business · Committees")
title(s, "Committees Need Chairs — That Means YOU \U0001F64B")
coms = ["\U0001F3C8 Sports / Games", "\U0001F3B6 Music", "\U0001F37D Food",
        "\U0001F9BA Safety", "\U0001F64F Services", "\U0001F4F8 Photography",
        "\U0001F49D Philanthropy & Fundraising", "\U0001F5D3 Planning", "\U0001F4BB Technology",
        "\U0001F490 Courtesy (led by the Hospitality Chair)", "\U0001F4DC Protocol (chaired by the Parliamentarian)"]
for i, cname in enumerate(coms):
    col, row = i % 3, i // 3
    x, y = 0.8 + col * 4.1, 1.7 + row * 0.85
    text(s, x, y, 3.9, 0.95, [(cname, {"size": 16, "bold": True, "color": GREEN, "space_after": 1}),
                               ("Chair: [sign up] · Members: [sign up]", {"size": 12, "color": MUTED})])
bullets(s, 0.7, 5.25, 6.2, [
    "Chairs appointed by the President, confirmed by member vote",
    "Meet monthly · report to the Vice President · 3-minute report each meeting",
    "Budget requests in writing to the Treasurer BEFORE spending"],
    head="How committees operate", size=13, h=2.0)
bullets(s, 7.2, 5.25, 5.4, [
    "Paper sheets at the welcome table (write your name under Chair or Member)",
    "Online: Committee Sign-Up form on hnfamilyreunion.com",
    "Right now: raise your hand when your committee is called"],
    head="Three ways to sign up", size=13, h=2.0, head_color=ORANGE)

# -------------------------------------------------------- 19 · NEXT LOCATION
s = slide()
kicker(s, "Family Business · Decade Planning")
title(s, "Where We Meet Next \U0001F5FA — Halfway Between Jackson, MS & Cleveland, OH")
text(s, 0.7, 1.7, W - 1.4, 0.6,
     "The drive is ~930 miles, so the halfway point lands around Nashville, TN / Bowling Green, KY.",
     size=15, color=MUTED)
bullets(s, 0.7, 2.4, 4.0, [
    "Montgomery Bell State Park — Burns, TN: 117 lakefront lodge rooms, 7,000+ sq ft event space",
    "Barren River Lake State Resort Park — Bowling Green, KY: banquet room seats 224–400"],
    head="\U0001F3AF At the halfway mark", size=13)
bullets(s, 4.85, 2.4, 4.0, [
    "Kentucky Dam Village State Resort Park — Gilbertsville, KY",
    "General Butler State Resort Park — off I-71, Louisville↔Cincinnati",
    "Brown County / Abe Martin Lodge — IN (indoor water park!)",
    "Hueston Woods Lodge — College Corner, OH"],
    head="\U0001F697 Along the route", size=13, head_color=ORANGE)
bullets(s, 9.0, 2.4, 3.9, [
    "Salt Fork State Park Lodge — Lore City, OH",
    "Punderson Manor Lodge — Newbury, OH",
    "The Lodge at Geneva-on-the-Lake — Lake Erie shore"],
    head="\U0001F3D9 Closer to Cleveland", size=13, head_color=RED)
footer(s, "The Planning Committee gets group-rate quotes from the top three picks; we choose by member vote. Decade planning = we rotate.")

# --------------------------------------------------------------- 20 · WEBSITE
s = slide()
kicker(s, "Family Business · Technology")
title(s, "Our New Home Online: hnfamilyreunion.com \U0001F4BB")
bullets(s, 0.7, 1.9, 4.0, ["Announcements & event schedule", "Family history & photo gallery",
                            "2024 slideshow & yearbook archive", "Memorial page"],
        head="Everything reunion", size=14)
bullets(s, 4.85, 1.9, 4.0, ["Registration & dues online", "T-shirt orders",
                             "Donations to the Land & Scholarship funds"],
        head="Payments & shop", size=14, head_color=ORANGE)
bullets(s, 9.0, 1.9, 3.9, ["Committee & chair sign-up", "Superlatives voting",
                            "Family tree submission form", "Contact the officers"],
        head="Get involved", size=14, head_color=RED)
footer(s, "The Technology Committee launches from starter code already written (in the family GitHub repo). Payments run through Zelle/Cash App and a secure processor — we never store card numbers.")

# ------------------------------------------------------------------ 21 · WIML
s = slide()
kicker(s, "Family Business · Our Land")
title(s, "Getting Our Family Land Back \U0001F3E1")
text(s, 0.7, 1.8, 6.2, 2.6, [
    ("We are partnering with Where Is My Land (whereismyland.com) — an organization that helps Black families research, document, and reclaim land that was taken from them.", {"size": 16, "space_after": 10}),
    ("Step 1: gather deeds, tax records & oral history\nStep 2: fund the RETAINER through reunion fundraising\nStep 3: family volunteers support the research", {"size": 15})])
bullets(s, 7.2, 1.8, 5.4, [
    "\U0001F3E1 A share of every fundraiser → the land-back fund (retainer + renovation & revitalization of our land)",
    "\U0001F6E0 A share → the OPERATING FUND (strictly operating costs: website hosting, printing, permits + seed money for future reunions)",
    "\U0001F393 A share → the scholarship fund",
    "Exact percentages set by member vote today — chair, entertain a motion!"],
    head="How fundraising is split", size=15, head_color=ORANGE)
footer(s, "If our claim succeeds, recovered land/proceeds are stewarded by the family foundation.")

# ---------------------------------------------------------- 22 · SCHOLARSHIPS
s = slide()
kicker(s, "Family Business · Education")
title(s, "The Harris-Nelson Scholarship Fund \U0001F393")
bullets(s, 0.7, 1.8, 6.0, [
    "Who: any descendant enrolled (or enrolling) in an accredited program — college, trade school, certifications",
    "How: short application + one-page essay on what family legacy means to you",
    "Selection: Philanthropy & Fundraising Committee reviews; officers ratify; awarded at the banquet"],
    head="A scholarship for descendants of our ancestors", size=15)
bullets(s, 7.0, 1.8, 5.6, [
    "A set share of all reunion fundraising (vote today)",
    "Dedicated fundraisers: raffle, fish fry, walk-a-thon, souvenirs",
    "Online giving — in honor or in memory of a loved one",
    "Tax-deductible once the 501(c)(3) is approved"],
    head="How we'll fund it", size=15, head_color=RED)
footer(s, "Goal: award our first scholarship(s) at the next reunion.")

# ---------------------------------------------------- 22b · HARDSHIP FUND
s = slide()
kicker(s, "Family Business · New Program")
title(s, "Helping Hands — the Family Hardship Fund \U0001F91D")
bullets(s, 0.7, 1.8, 4.0, [
    "Any descendant, spouse/partner, or child (minors via a guardian) — or apply on behalf of a relative",
    "Hardship within the past 12 months: medical, job loss, housing/disaster, funeral, caregiving, education emergency",
    "One award per household per cycle; dues status does NOT affect eligibility"],
    head="Who can apply", size=13)
bullets(s, 4.85, 1.8, 4.0, [
    "Confidential application → Philanthropy & Fundraising Committee scores it",
    "Rubric: severity & urgency 40% · impact 30% · clarity & completeness 20% · discretion 10%",
    "3+ independent scorers; household conflicts recuse; officers ratify",
    "Recipient may stay ANONYMOUS — always their choice"],
    head="How the recipient is chosen", size=13, head_color=ORANGE)
bullets(s, 9.0, 1.8, 3.9, [
    "Applications open January 1 of the reunion year",
    "Deadline: 30 days before the reunion",
    "Committee scores by 14 days out; officers ratify by 7",
    "Award presented at Sunday brunch"],
    head="Deadlines each cycle", size=13, head_color=RED)
footer(s, "Apply at hnfamilyreunion.com (Hardship Fund page) or on paper from any officer. Award amount set by the Treasurer's budget each cycle — funded from the fundraising split.")

# ----------------------------------------------------------------- 23 · 501c3
s = slide()
kicker(s, "Family Business · The Foundation")
title(s, "Starting Our Family Foundation — a 501(c)(3) \U0001F3DB")
bullets(s, 0.7, 1.8, 6.0, [
    "Month 1: choose the name, elect the founding board, draft bylaws & conflict-of-interest policy",
    "Month 1–2: file Articles of Incorporation; get an EIN (free, IRS.gov)",
    "Month 2–3: file IRS Form 1023-EZ (if eligible) or Form 1023",
    "Month 3+: state charity registration, bank account, annual budget"],
    head="The process (next few months)", size=15)
bullets(s, 7.0, 1.8, 5.6, [
    "Board Chair · Vice Chair · Board Secretary · Board Treasurer",
    "3–5 at-large directors — one from each generation if we can",
    "Its own body with its own votes, separate from reunion officers",
    "Board meets quarterly; minutes published to the family"],
    head="Foundation board (to be elected)", size=15, head_color=ORANGE)
footer(s, "Nominations open today — see the Secretary. Already raised their hands on the questionnaire: Toni Nelson & Dennis Derrick Jr. We ratify the founding board by motion and vote.")

# ----------------------------------------------------------- 24 · FAMILY TREE
s = slide()
kicker(s, "All Weekend Long")
title(s, "The Family Tree Project — Add Yourself! \U0001F333")
bullets(s, 0.7, 1.8, 6.0, [
    "Stop by the Family Tree table (or scan the QR code)",
    "Fill out the short form — name, birthday, parents, branch, children",
    "The Historian (Shone) adds you to the Ancestry tree",
    "A large family tree poster gets printed with EVERYONE listed"],
    head="Four easy steps", size=15)
bullets(s, 7.0, 1.8, 5.6, [
    "Online: the Family Tree Google Form, or the matching form on hnfamilyreunion.com",
    "Paper: forms at the Family Tree table all weekend",
    "Bring old photos, obituaries, bibles, deeds — we scan them on the spot (you keep the originals)"],
    head="Ways to submit", size=15, head_color=ORANGE)

# ------------------------------------------------------------------ 25 · LINKS
s = slide()
kicker(s, "Keep These Handy"); title(s, "All the Links \U0001F517")
links = [("Family website", "hnfamilyreunion.com (launching soon)"),
         ("Family email", "harrisnelsonfamilyreunion@gmail.com"),
         ("Zelle (dues & shirts)", "Joseph Nelson — harrisnelsonfamilyreunion@gmail.com"),
         ("Cash App", "Miesha Wilson — $mieshanulife06"),
         ("Yearbook post", "x.com/tendin2"),
         ("History & yearbook doc", "Google Docs (link on the website)"),
         ("Photos & stories album", "canva.link/fqxj2trjluu9fl5"),
         ("Family tree form", "Google Form (QR at the Family Tree table)"),
         ("Land-back partner", "whereismyland.com")]
tbl = s.shapes.add_table(len(links) + 1, 2, Inches(0.8), Inches(1.8),
                         Inches(11.7), Inches(4.8)).table
tbl.columns[0].width = Inches(3.7); tbl.columns[1].width = Inches(8.0)
for j, htxt in enumerate(["What", "Where"]):
    c = tbl.cell(0, j); c.text = htxt
    c.fill.solid(); c.fill.fore_color.rgb = GREEN
    r = c.text_frame.paragraphs[0].runs[0].font
    r.size = Pt(12); r.bold = True; r.color.rgb = WHITE
for i, row in enumerate(links, 1):
    for j, v in enumerate(row):
        c = tbl.cell(i, j); c.text = v
        f = c.text_frame.paragraphs[0].runs[0].font
        f.size = Pt(13); f.color.rgb = INK; f.bold = (j == 0)

# ---------------------------------------------------------------- 26 · CLOSING
s = slide()
text(s, 0.7, 2.1, W - 1.4, 1.9, "We Are the Harvest\nof Their Planting", size=46,
     color=GREEN, bold=True, align=PP_ALIGN.CENTER)
text(s, 0.7, 4.2, W - 1.4, 1.0,
     "Sign a form. Chair a committee. Add your branch to the tree.\nSee you at the next reunion — bigger, better, and on our own land.",
     size=18, color=MUTED, align=PP_ALIGN.CENTER)
text(s, 0.7, 5.6, W - 1.4, 0.5,
     "Harris-Nelson Family Reunion · hnfamilyreunion.com · harrisnelsonfamilyreunion@gmail.com",
     size=13, color=MUTED, align=PP_ALIGN.CENTER, font="Verdana")

out = os.path.join(HERE, "Harris-Nelson-Reunion-2026.pptx")
prs.save(out)
print("saved", out, f"({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
